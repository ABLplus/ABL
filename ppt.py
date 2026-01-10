from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(10, 25, 47)      # Deep Navy
TEXT_COLOR = RGBColor(255, 255, 255) # White
ACCENT_COLOR = RGBColor(255, 215, 0) # Gold
SEC_ACCENT = RGBColor(100, 255, 218) # Teal/Cyan

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, font_size=44, color=ACCENT_COLOR):
    title_shape = slide.shapes.title
    title_shape.text = text
    title_shape.text_frame.paragraphs[0].font.color.rgb = color
    title_shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.top = Inches(0.5)

def add_text_box(slide, text, left, top, width, height, font_size=20, color=TEXT_COLOR, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- SLIDE 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    set_slide_background(slide)
    
    # Title Text
    tf = add_text_box(slide, "CLEAR PRELIMS 2026", Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1), 60, ACCENT_COLOR, True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    tf = add_text_box(slide, "A Community Session by ABLE IAS & ABL+", Inches(0), Inches(3.5), SLIDE_WIDTH, Inches(1), 32, TEXT_COLOR)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Details
    tf = add_text_box(slide, "Tonight, 8 PM | For Serious Aspirants Only\nSpeaker: [Your Name]", Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(1), 20, SEC_ACCENT)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- SLIDE 2: Reality ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Reality Aspirants Face")
    
    # Left Content
    content = "Prelims feels uncertain because:\n\n• Difficulty fluctuates\n• Weightage varies\n• Options are tricky\n• CA linkages complicate\n• Everything feels scattered"
    add_text_box(slide, content, Inches(1), Inches(2), Inches(5), Inches(4), 24)
    
    # Right Visual (Box)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.5), Inches(5), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(40, 50, 70)
    shape.line.color.rgb = SEC_ACCENT
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "But uncertainty is\nemotional — not structural."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True

    # --- SLIDE 3: The Structure (Visual Pyramid) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Paper Has a Structure")
    
    # Draw Layered Blocks
    # Peripheral (Top)
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2), Inches(5.33), Inches(1))
    s1.fill.fore_color.rgb = RGBColor(100, 100, 100)
    s1.text_frame.text = "Peripheral Zone (≈ 30–35 Qs)"
    
    # Derivative (Middle)
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.1), Inches(6.33), Inches(1.2))
    s2.fill.fore_color.rgb = RGBColor(200, 100, 50) # Orange
    s2.text_frame.text = "Core Derivatives (≈ 30–35 Qs)"
    
    # Core (Bottom)
    s3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.4), Inches(7.33), Inches(1.5))
    s3.fill.fore_color.rgb = RGBColor(200, 50, 50) # Red
    s3.text_frame.text = "Core of Core (≈ 35 Qs)\n(Base Foundation)"

    # Bottom Caption
    tf = add_text_box(slide, "Understanding this structure reduces 70% of anxiety.", Inches(1), Inches(6.2), Inches(11), Inches(1), 24, SEC_ACCENT)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- SLIDE 4: Core of Core ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Core of Core (≈ 35 Qs)")
    
    # Content
    content = "Fundamentals of:\n• Polity\n• Economy\n• History\n• Geography\n• Basics of Environment & S&T"
    add_text_box(slide, content, Inches(1.5), Inches(2), Inches(5), Inches(4), 24)
    
    # Target Box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2.5), Inches(4.5), Inches(2))
    shape.fill.fore_color.rgb = RGBColor(0, 100, 0) # Green
    tf = shape.text_frame
    tf.text = "TARGET:\nGet 30/35 correct\n≈ 56+ marks"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- SLIDE 6: Non-Negotiable Truth ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    # Center Big Text
    tf = add_text_box(slide, "THE NON-NEGOTIABLE TRUTH", Inches(1), Inches(1), Inches(11.33), Inches(1), 40, ACCENT_COLOR, True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    tf = add_text_box(slide, "You must solve enough questions correctly.\nYou must convert knowledge → marks.\n\nThere is no substitute for correct attempts.", Inches(2), Inches(3), Inches(9.33), Inches(3), 32, TEXT_COLOR)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- SLIDE 8: Practice Equation ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Practicing = Attempting")
    
    # Equation Visual
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(2.5), Inches(2.5), Inches(2.5))
    shape.text_frame.text = "Thinking"
    
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(3.5), Inches(1), Inches(0.5))
    
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(2.5), Inches(2.5), Inches(2.5))
    shape.text_frame.text = "Risk"
    
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.3), Inches(3.5), Inches(1), Inches(0.5))
    
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.4), Inches(2.5), Inches(2.5), Inches(2.5))
    shape.text_frame.text = "Instinct"

    tf = add_text_box(slide, "Every question reveals your gaps.", Inches(0), Inches(6), SLIDE_WIDTH, Inches(1), 24, SEC_ACCENT)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- SLIDE 12: Monitoring Graph ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Monitoring & Accountability")
    
    # Draw simple graph axis
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(5.5), Inches(8), Inches(0.05)) # X Axis
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(0.05), Inches(3.5)) # Y Axis
    
    # Draw rising line
    line = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2), Inches(5.5), Inches(8), Inches(0.2))
    line.rotation = -30
    line.fill.fore_color.rgb = SEC_ACCENT
    
    add_text_box(slide, "Reliability → Confidence → Performance", Inches(3), Inches(6), Inches(7), Inches(1), 24, ACCENT_COLOR, True)

    # --- SLIDE 15: CTA ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    tf = add_text_box(slide, "Join ABL+", Inches(0), Inches(2), SLIDE_WIDTH, Inches(1), 50, ACCENT_COLOR, True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content = "Strengthen preparation with structured attempts.\n\nTry for Free → www.abc.com\nTelegram → t.me/abc"
    tf = add_text_box(slide, content, Inches(0), Inches(3.5), SLIDE_WIDTH, Inches(3), 28, TEXT_COLOR)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save
    prs.save('Clear_Prelims_2026.pptx')
    print("Presentation generated successfully with visuals!")

if __name__ == "__main__":
    create_presentation()